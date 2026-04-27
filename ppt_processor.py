"""
PowerPoint Processing Module
Handles PPT file operations with GPT-5.4 enhanced desensitization
"""

import os
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import logging
from dataclasses import dataclass

logger = logging.getLogger(__name__)

@dataclass
class PPTElement:
    """Represents a PPT element that can contain text"""
    element_type: str
    text: str
    location: str
    slide_number: int
    shape_id: Optional[str] = None
    is_chart: bool = False
    is_table: bool = False
    is_title: bool = False

class PowerPointProcessor:
    """
    Advanced PowerPoint processor with GPT-5.4 enhanced desensitization
    """
    
    def __init__(self, desensitizer):
        self.desensitizer = desensitizer
        self.processed_elements = []
        self.statistics = {
            'slides_processed': 0,
            'text_elements_processed': 0,
            'charts_processed': 0,
            'tables_processed': 0,
            'redactions_applied': 0,
            'preserved_elements': 0
        }
    
    def process_presentation(self, input_path: str, output_path: str) -> bool:
        """
        Process a PowerPoint presentation for desensitization
        """
        try:
            logger.info(f"Processing presentation: {input_path}")
            
            # Load presentation
            prs = Presentation(input_path)
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                self._process_slide(slide, slide_idx + 1)
                self.statistics['slides_processed'] += 1
            
            # Save processed presentation
            prs.save(output_path)
            
            # Generate audit report
            self._generate_audit_report(output_path)
            
            logger.info(f"Presentation processed successfully: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to process presentation: {e}")
            return False
    
    def _process_slide(self, slide, slide_number: int):
        """Process individual slide"""
        logger.debug(f"Processing slide {slide_number}")
        
        # Process slide elements
        for shape in slide.shapes:
            self._process_shape(shape, slide_number)
    
    def _process_shape(self, shape, slide_number: int):
        """Process individual shape on slide"""
        try:
            # Handle different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                self._process_text_shape(shape, slide_number)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self._process_table_shape(shape, slide_number)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                self._process_chart_shape(shape, slide_number)
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                self._process_text_shape(shape, slide_number)
                
        except Exception as e:
            logger.warning(f"Error processing shape on slide {slide_number}: {e}")
    
    def _process_text_shape(self, shape, slide_number: int):
        """Process text-containing shapes"""
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return
        
        is_title = self._is_title_shape(shape)
        original_text = shape.text_frame.text
        
        # Create context for GPT analysis
        context = self._build_context(shape, slide_number, is_title)
        
        # Apply desensitization
        processed_text = self.desensitizer.desensitize_text(original_text, context)
        
        # Update text if changed
        if processed_text != original_text:
            self._update_text_frame(shape.text_frame, processed_text)
            self.statistics['redactions_applied'] += 1
            
            # Record for audit
            element = PPTElement(
                element_type="text_shape",
                text=original_text,
                location=f"slide_{slide_number}_shape_{shape.shape_id}",
                slide_number=slide_number,
                shape_id=str(shape.shape_id),
                is_title=is_title
            )
            self.processed_elements.append(element)
        else:
            self.statistics['preserved_elements'] += 1
        
        self.statistics['text_elements_processed'] += 1
    
    def _process_table_shape(self, shape, slide_number: int):
        """Process table shapes"""
        if not hasattr(shape, 'table') or not shape.table:
            return
        
        table = shape.table
        self.statistics['tables_processed'] += 1
        
        # Process each cell in the table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame and cell.text_frame.text:
                    original_text = cell.text_frame.text
                    
                    # Build context for table cell
                    context = f"table_cell_row_{row_idx}_col_{col_idx}_slide_{slide_number}"
                    
                    # Apply desensitization
                    processed_text = self.desensitizer.desensitize_text(original_text, context)
                    
                    # Update cell text if changed
                    if processed_text != original_text:
                        self._update_text_frame(cell.text_frame, processed_text)
                        self.statistics['redactions_applied'] += 1
    
    def _process_chart_shape(self, shape, slide_number: int):
        """Process chart shapes"""
        self.statistics['charts_processed'] += 1
        
        # Chart processing would require more complex XML manipulation
        # For now, we'll handle chart titles and labels
        if hasattr(shape, 'chart') and shape.chart:
            chart = shape.chart
            
            # Process chart title
            if chart.has_title and chart.chart_title:
                original_title = chart.chart_title.text_frame.text
                context = f"chart_title_slide_{slide_number}"
                processed_title = self.desensitizer.desensitize_text(original_title, context)
                
                if processed_title != original_title:
                    chart.chart_title.text_frame.text = processed_title
                    self.statistics['redactions_applied'] += 1
    
    def _is_title_shape(self, shape) -> bool:
        """Determine if shape is a title"""
        # Check if it's in title placeholder
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            placeholder_type = shape.placeholder_format.type
            return placeholder_type in [1, 3]  # Title or subtitle placeholder
        
        # Check text properties (size, position)
        if hasattr(shape, 'top') and hasattr(shape, 'left'):
            # Simple heuristic: if it's at the top and large, likely a title
            return shape.top < 100
        
        return False
    
    def _build_context(self, shape, slide_number: int, is_title: bool) -> str:
        """Build context for GPT analysis"""
        context_parts = [f"slide_{slide_number}"]
        
        if is_title:
            context_parts.append("title")
        
        # Add position context
        if hasattr(shape, 'top') and hasattr(shape, 'left'):
            if shape.top < 100:
                context_parts.append("header")
            elif shape.top > 400:
                context_parts.append("footer")
        
        return "_".join(context_parts)
    
    def _update_text_frame(self, text_frame, new_text: str):
        """Update text frame content"""
        # Clear existing text
        text_frame.clear()
        
        # Add new text
        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        p.text = new_text
    
    def _generate_audit_report(self, output_path: str):
        """Generate audit report of desensitization process"""
        report_path = output_path.replace('.pptx', '_audit_report.json')
        
        audit_report = {
            'statistics': self.statistics,
            'processed_elements': [
                {
                    'element_type': elem.element_type,
                    'original_text': elem.text,
                    'location': elem.location,
                    'slide_number': elem.slide_number,
                    'is_title': elem.is_title
                }
                for elem in self.processed_elements
            ],
            'timestamp': self._get_timestamp(),
            'tool_version': '1.0.0'
        }
        
        try:
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(audit_report, f, ensure_ascii=False, indent=2)
            logger.info(f"Audit report generated: {report_path}")
        except Exception as e:
            logger.error(f"Failed to generate audit report: {e}")
    
    def _get_timestamp(self) -> str:
        """Get current timestamp"""
        from datetime import datetime
        return datetime.now().isoformat()
    
    def get_processing_summary(self) -> Dict:
        """Get processing summary"""
        return {
            'total_elements_processed': len(self.processed_elements),
            'statistics': self.statistics,
            'redaction_rate': self.statistics['redactions_applied'] / max(self.statistics['text_elements_processed'], 1),
            'preservation_rate': self.statistics['preserved_elements'] / max(self.statistics['text_elements_processed'], 1)
        }

# Advanced PPT XML processing for deeper desensitization
class PPTXMLProcessor:
    """
    Advanced XML processor for PowerPoint files
    Handles embedded XML content for comprehensive desensitization
    """
    
    def __init__(self, desensitizer):
        self.desensitizer = desensitizer
    
    def extract_xml_content(self, pptx_path: str) -> Dict[str, str]:
        """Extract XML content from PPTX file"""
        xml_contents = {}
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                # Extract slide XML files
                for filename in zip_ref.namelist():
                    if filename.startswith('ppt/slides/slide') and filename.endswith('.xml'):
                        xml_content = zip_ref.read(filename).decode('utf-8')
                        xml_contents[filename] = xml_content
        
        except Exception as e:
            logger.error(f"Failed to extract XML content: {e}")
        
        return xml_contents
    
    def process_xml_content(self, xml_content: str, context: str = "") -> str:
        """Process XML content for desensitization"""
        try:
            # Parse XML
            root = ET.fromstring(xml_content)
            
            # Find all text elements
            text_elements = root.findall(".//a:t", namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            
            # Process each text element
            for elem in text_elements:
                if elem.text:
                    processed_text = self.desensitizer.desensitize_text(elem.text, context)
                    if processed_text != elem.text:
                        elem.text = processed_text
            
            # Return processed XML
            return ET.tostring(root, encoding='unicode')
            
        except Exception as e:
            logger.error(f"Failed to process XML content: {e}")
            return xml_content

# Usage example
if __name__ == "__main__":
    from ppt_desensitizer import GPTEnhancedDesensitizer
    
    # Initialize desensitizer
    desensitizer = GPTEnhancedDesensitizer('desensitization_config.json')
    
    # Create processor
    processor = PowerPointProcessor(desensitizer)
    
    # Example usage (would need actual PPT file)
    # success = processor.process_presentation('input.pptx', 'output.pptx')
    
    # Print statistics
    summary = processor.get_processing_summary()
    print("Processing Summary:")
    for key, value in summary.items():
        print(f"{key}: {value}")